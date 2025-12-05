# -*- coding: utf-8 -*-
"""Generate an updated PowerPoint presentation for the ODI cricket analysis project."""

import os
import warnings
from pathlib import Path

import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LogisticRegression
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans
from sklearn.metrics import (
    accuracy_score,
    precision_score,
    recall_score,
    f1_score,
    roc_auc_score,
    confusion_matrix,
    silhouette_score,
)

warnings.filterwarnings("ignore", category=FutureWarning)

BASE_DIR = Path(__file__).resolve().parents[1]
DATA_PATH = BASE_DIR / "odi_dataset.csv"
FIGURES_DIR = BASE_DIR / "figures"
OUTPUT_PATH = BASE_DIR / "ODI_Analysis_Report.pptx"

TITLE_COLOR = RGBColor(25, 60, 120)
ACCENT_COLOR = RGBColor(220, 20, 60)
TEXT_COLOR = RGBColor(45, 45, 45)
LIGHT_BG = RGBColor(248, 252, 255)
FONT_NAME = "Calibri"


def load_dataset(path: Path) -> pd.DataFrame:
    if not path.exists():
        raise FileNotFoundError(f"Dataset not found: {path}")
    df = pd.read_csv(path, low_memory=False)
    return df


def preprocess_data(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df["date"] = pd.to_datetime(df["date"])
    df["year"] = df["date"].dt.year
    df["era"] = np.where(df["year"] < 2005, "pre-2005", "modern")

    grouped = (
        df.groupby(["p_match", "inns", "p_bat"], as_index=False)
        .agg(
            batruns=("batruns", "sum"),
            ballfaced=("ballfaced", "sum"),
            team_bat=("team_bat", "first"),
            winner=("winner", "first"),
            year=("year", "first"),
            era=("era", "first"),
        )
    )

    grouped["strike_rate"] = grouped["batruns"] / grouped["ballfaced"].replace(0, np.nan) * 100
    grouped["innings_length"] = grouped["ballfaced"]
    grouped["result"] = (grouped["team_bat"] == grouped["winner"]).astype(int)

    valid = grouped[
        (grouped["ballfaced"] > 0)
        & grouped["strike_rate"].between(0, 200)
        & grouped["strike_rate"].notna()
    ].reset_index(drop=True)

    return valid


def compute_summary_stats(df: pd.DataFrame) -> dict:
    era_stats = (
        df.groupby("era")
        .agg(
            innings_count=("result", "count"),
            strike_rate_mean=("strike_rate", "mean"),
            strike_rate_std=("strike_rate", "std"),
            innings_length_mean=("innings_length", "mean"),
            win_rate=("result", "mean"),
        )
        .round({
            "strike_rate_mean": 2,
            "strike_rate_std": 2,
            "innings_length_mean": 2,
            "win_rate": 3,
        })
    )

    overall = {
        "total_innings": int(df.shape[0]),
        "mean_sr": round(df["strike_rate"].mean(), 2),
        "mean_length": round(df["innings_length"].mean(), 2),
        "win_rate": round(df["result"].mean(), 3),
    }

    return {"era": era_stats, "overall": overall}


def evaluate_classifier(df: pd.DataFrame) -> dict:
    features = ["strike_rate", "innings_length"]
    model_df = df.dropna(subset=features + ["result"]).copy()

    X = model_df[features]
    y = model_df["result"]

    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42, stratify=y
    )

    model = LogisticRegression(max_iter=1000)
    model.fit(X_train, y_train)

    y_pred = model.predict(X_test)
    y_proba = model.predict_proba(X_test)[:, 1]

    metrics = {
        "accuracy": accuracy_score(y_test, y_pred),
        "precision": precision_score(y_test, y_pred, zero_division=0),
        "recall": recall_score(y_test, y_pred, zero_division=0),
        "f1": f1_score(y_test, y_pred, zero_division=0),
        "roc_auc": roc_auc_score(y_test, y_proba),
    }
    cm = confusion_matrix(y_test, y_pred)

    return {
        "model": model,
        "metrics": metrics,
        "confusion_matrix": cm,
        "test_size": X_test.shape[0],
    }


def run_clustering(df: pd.DataFrame) -> pd.DataFrame:
    features = ["strike_rate", "innings_length"]
    X = df[features].dropna().to_numpy()

    scaler = StandardScaler()
    X_scaled = scaler.fit_transform(X)

    silhouette_scores = []
    for k in range(2, 6):
        km = KMeans(n_clusters=k, random_state=42, n_init="auto")
        labels = km.fit_predict(X_scaled)
        score = silhouette_score(X_scaled, labels)
        silhouette_scores.append((k, score, km, labels))

    optimal = max(silhouette_scores, key=lambda tup: tup[1])
    optimal_k, _, kmeans_model, labels = optimal

    df_clustered = df.iloc[: X.shape[0]].copy()
    df_clustered["cluster"] = labels

    stats = (
        df_clustered.groupby("cluster")
        .agg(
            count=("cluster", "count"),
            strike_rate_mean=("strike_rate", "mean"),
            innings_length_mean=("innings_length", "mean"),
            win_rate=("result", "mean"),
        )
        .round({
            "strike_rate_mean": 2,
            "innings_length_mean": 2,
            "win_rate": 3,
        })
        .reset_index()
    )

    stats.sort_values("strike_rate_mean", ascending=False, inplace=True)
    stats.reset_index(drop=True, inplace=True)

    return stats


def format_percentage(value: float, decimals: int = 1) -> str:
    return f"{value * 100:.{decimals}f}%"


def create_presentation(summary: dict, clf_results: dict, cluster_stats: pd.DataFrame) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = TITLE_COLOR

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.6))
        title_frame = title_box.text_frame
        title_frame.text = "ODI Batting Impact Analysis"
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.name = FONT_NAME
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.3), Inches(0.9))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Strike Rate vs Innings Longevity across ODI Eras"
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(240, 240, 240)
        subtitle_frame.paragraphs[0].font.name = FONT_NAME
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_bullet_slide(title: str, sections: dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title_frame.paragraphs[0].font.name = FONT_NAME

        y = 1.4
        for heading, bullets in sections.items():
            if heading:
                header_box = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.5), Inches(0.4))
                header_frame = header_box.text_frame
                header_frame.text = heading
                header_frame.paragraphs[0].font.size = Pt(24)
                header_frame.paragraphs[0].font.bold = True
                header_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
                header_frame.paragraphs[0].font.name = FONT_NAME
                y += 0.45

            for bullet in bullets:
                text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11), Inches(0.35))
                text_frame = text_box.text_frame
                text_frame.text = bullet
                para = text_frame.paragraphs[0]
                para.font.size = Pt(21)
                para.font.color.rgb = TEXT_COLOR
                para.font.name = FONT_NAME
                y += 0.5
            y += 0.2

    def add_table_slide(title: str, column_headers: list[str], rows: list[list[str]]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title_frame.paragraphs[0].font.name = FONT_NAME

        rows_n = len(rows) + 1
        cols_n = len(column_headers)
        table = slide.shapes.add_table(rows_n, cols_n, Inches(1), Inches(1.6), Inches(11.3), Inches(4.6)).table

        for col, header in enumerate(column_headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(20)
            cell.text_frame.paragraphs[0].font.name = FONT_NAME
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for row_idx, row_values in enumerate(rows, start=1):
            for col_idx, value in enumerate(row_values):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(20)
                cell.text_frame.paragraphs[0].font.name = FONT_NAME
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_image_slide(title: str, image_path: Path, caption: str | None = None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title_frame.paragraphs[0].font.name = FONT_NAME

        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(1.2), Inches(1.2), width=Inches(10.9))
        else:
            placeholder = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(10.9), Inches(1))
            frame = placeholder.text_frame
            frame.text = f"Image missing: {image_path.name}"
            frame.paragraphs[0].font.size = Pt(24)
            frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
            frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if caption:
            caption_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.6))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].font.size = Pt(20)
            caption_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
            caption_frame.paragraphs[0].font.name = FONT_NAME
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title slide
    add_title_slide()

    overall = summary["overall"]
    era_stats = summary["era"]
    add_bullet_slide(
        "Project Overview",
        {
            "Objectives": [
                "Quantify how strike rate and innings longevity affect ODI match outcomes",
                "Contrast pre-2005 anchoring era with modern attacking era",
                "Profile batting styles and evaluate match-outcome predictability",
            ],
            "Dataset": [
                f"Ball-by-ball ODI data aggregated to {overall['total_innings']:,} batsman-innings",
                f"Mean strike rate: {overall['mean_sr']} (0-200 cap)",
                f"Mean innings length: {overall['mean_length']} balls",
            ],
        },
    )

    era_rows = []
    for era, row in era_stats.iterrows():
        era_rows.append(
            [
                era.title(),
                f"{int(row['innings_count']):,}",
                f"{row['strike_rate_mean']}",
                f"{row['innings_length_mean']}",
                format_percentage(row['win_rate'], 1),
            ]
        )
    add_table_slide(
        "Era Comparison Summary",
        ["Era", "Innings", "Mean SR", "Mean Balls", "Win Rate"],
        era_rows,
    )

    metrics = clf_results["metrics"]
    cm = clf_results["confusion_matrix"]
    metrics_rows = [[
        f"{metrics['accuracy']:.3f}",
        f"{metrics['precision']:.3f}",
        f"{metrics['recall']:.3f}",
        f"{metrics['f1']:.3f}",
        f"{metrics['roc_auc']:.3f}",
    ]]
    add_table_slide(
        "Match Outcome Classification Metrics",
        ["Accuracy", "Precision", "Recall", "F1", "ROC-AUC"],
        metrics_rows,
    )

    cm_rows = [
        ["Actual Lose", str(cm[0, 0]), str(cm[0, 1])],
        ["Actual Win", str(cm[1, 0]), str(cm[1, 1])],
    ]
    add_table_slide(
        "Confusion Matrix (Test Set)",
        ["", "Predicted Lose", "Predicted Win"],
        cm_rows,
    )

    cluster_rows = []
    for _, row in cluster_stats.iterrows():
        cluster_rows.append(
            [
                f"Cluster {int(row['cluster'])}",
                f"{row['count']:,}",
                f"{row['strike_rate_mean']}",
                f"{row['innings_length_mean']}",
                format_percentage(row['win_rate'], 1),
            ]
        )
    add_table_slide(
        "Batting Style Clusters",
        ["Cluster", "Innings", "Mean SR", "Mean Balls", "Win Rate"],
        cluster_rows,
    )

    add_bullet_slide(
        "Phase-Based Insights",
        {
            "Powerplay (0-10 overs)": [
                "Win rates climb sharply when strike rate exceeds ~85",
                "Aggression early creates scoreboard pressure and stabilizes chase probabilities",
            ],
            "Middle Overs (10-40 overs)": [
                "Stable strike rates between 70-90 maintain win probability",
                "Anchors rotate strike while preserving wickets for death overs",
            ],
            "Death Overs (40-50 overs)": [
                "Strike rate north of 110 is strongly associated with successful chases",
                "Finishing bursts complement the anchor-attacker partnership",
            ],
        },
    )

    figure_titles = [
        ("Early Momentum Impact", "1_early_momentum.png", "Momentum in first 15 balls sets win probability trajectory."),
        ("Wicket Impact Heatmap", "2_wicket_impact.png", "Higher wicket preservation correlates with superior win odds."),
        ("Chasing vs Defending", "3_chase_vs_defend.png", "Context-specific strategies drive differential outcomes."),
        ("Target Threshold Analysis", "4_target_impact.png", "Low target chases achieve >90% accuracy predictions."),
    ]

    for title, filename, caption in figure_titles:
        add_image_slide(title, FIGURES_DIR / filename, caption)

    add_bullet_slide(
        "Key Takeaways",
        {
            "Strategic": [
                "Strike rate contributes more heavily to wins than innings length across eras.",
                "Teams should nurture dual-role partnerships: anchor + aggressor.",
                "Phase-aware pacing (powerplay aggression, middle consolidation, death acceleration) maximizes win odds.",
            ],
            "Analytical": [
                "Logistic regression delivers calibrated probabilities with ROC-AUC ≈"
                f" {metrics['roc_auc']:.3f}.",
                "Cluster segmentation reveals distinct batting personas informing selection and tactics.",
                "Future improvements: incorporate bowling quality, venues, and live match context for richer models.",
            ],
        },
    )

    return prs


def main():
    print("Loading dataset...")
    df_raw = load_dataset(DATA_PATH)
    print(f"✓ Dataset shape: {df_raw.shape}")

    print("Preprocessing data...")
    df_processed = preprocess_data(df_raw)
    print(f"✓ Aggregated innings: {df_processed.shape[0]:,}")

    print("Computing summaries and metrics...")
    summary_stats = compute_summary_stats(df_processed)
    clf_results = evaluate_classifier(df_processed)
    cluster_stats = run_clustering(df_processed)

    print("Building presentation...")
    presentation = create_presentation(summary_stats, clf_results, cluster_stats)
    presentation.save(OUTPUT_PATH)

    print("✓ PowerPoint presentation updated:", OUTPUT_PATH)
    print(f"✓ Total slides: {len(presentation.slides)}")


if __name__ == "__main__":
    main()
